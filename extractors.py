"""Document text extraction for the ChatGB10 router. Lazy imports so a missing
library only disables that one file type; the router still starts."""
import io

def extract_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    out = [p.text for p in doc.paragraphs if p.text.strip()]
    for t in doc.tables:
        for row in t.rows:
            out.append("\t".join(c.text.strip() for c in row.cells))
    return "\n".join(out)

def extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in vals):
                out.append("\t".join(vals))
    return "\n".join(out)

def extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs)
                    if txt.strip():
                        out.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    out.append("\t".join(c.text for c in row.cells))
    return "\n".join(out)

def extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    out = []
    for i, page in enumerate(reader.pages, 1):
        txt = (page.extract_text() or "").strip()
        if txt:
            out.append(f"# Page {i}\n{txt}")
    return "\n".join(out) or "(No selectable text found — this PDF may be a scan/image. OCR is not enabled.)"

def extract_text(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")

EXTRACTORS = {
    ".docx": extract_docx, ".xlsx": extract_xlsx, ".xlsm": extract_xlsx,
    ".pptx": extract_pptx, ".pdf": extract_pdf,
    ".txt": extract_text, ".md": extract_text, ".csv": extract_text,
    ".json": extract_text, ".log": extract_text, ".py": extract_text,
}
